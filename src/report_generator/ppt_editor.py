"""PowerPoint editing functionality."""
import logging
import tempfile
from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from .utils import atomic_write, generate_timestamp, safe_path

logger = logging.getLogger(__name__)

class PPTEditor:
    """Handle PowerPoint file editing operations."""
    
    def __init__(self, input_path: Path):
        self.input_path = input_path
        self.presentation = None
    
    def load_presentation(self) -> None:
        """Load PowerPoint presentation."""
        try:
            self.presentation = Presentation(self.input_path)
            logger.info(f"Loaded presentation with {len(self.presentation.slides)} slides")
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise
    
    def replace_text_in_slide(self, slide_number: int, replacements: Dict[str, Any]) -> None:
        """Replace text placeholders in specified slide."""
        if not self.presentation:
            raise ValueError("Presentation not loaded")
        
        if slide_number > len(self.presentation.slides):
            raise ValueError(f"Slide {slide_number} does not exist")
        
        slide = self.presentation.slides[slide_number - 1]  # 0-indexed
        replacements_made = 0
        
        # Process all shapes in slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, value in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))
                                replacements_made += 1
        
        logger.info(f"Made {replacements_made} replacements in slide {slide_number}")
    
    def save_report(self, output_dir: Path) -> Path:
        """Save modified presentation with timestamped filename."""
        if not self.presentation:
            raise ValueError("Presentation not loaded")
        
        timestamp = generate_timestamp()
        filename = f"report_generated_{timestamp}.pptx"
        output_path = safe_path(output_dir, filename)
        
        # Save to temporary file first
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            self.presentation.save(tmp_file.name)
            atomic_write(output_path, Path(tmp_file.name))
            Path(tmp_file.name).unlink()  # Clean up temp file
        
        logger.info(f"Report saved to {output_path}")
        return output_path