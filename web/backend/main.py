"""FastAPI backend for PPT report generator."""
import json
import tempfile
from pathlib import Path
from typing import Dict
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
import sys
sys.path.append(str(Path(__file__).parent.parent.parent))
from src.report_generator.ppt_editor import PPTEditor


app = FastAPI(title="Automated Executive Report Builder", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:8000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class ReplacementRequest(BaseModel):
    replacements: Dict[str, str]
    slide_number: int = 2

@app.get("/api/config")
async def get_default_config():
    """Get default configuration and slide text."""
    # Get project root directory
    project_root = Path(__file__).parent.parent.parent
    config_path = project_root / "report_config.json"
    
    try:
        with open(config_path) as f:
            config = json.load(f)
        
        # Get slide text
        ppt_path = project_root / "input_PPT" / "SMBC.pptx"
        if ppt_path.exists():
            pres = Presentation(ppt_path)
            slide = pres.slides[config["slide_number"] - 1]
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    slide_text += shape.text_frame.text.strip() + "\n\n"
            
            return {
                "config": config,
                "slide_text": slide_text.strip()
            }
        else:
            return {"config": config, "slide_text": "PPT file not found"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate")
async def generate_report(request: ReplacementRequest):
    """Generate PPT report with replacements."""
    try:
        # Get project root directory
        project_root = Path(__file__).parent.parent.parent
        ppt_path = project_root / "input_PPT" / "SMBC.pptx"
        if not ppt_path.exists():
            raise HTTPException(status_code=404, detail="PPT template not found")
        
        # Use project output directory instead of temp
        output_dir = project_root / "output"
        output_dir.mkdir(exist_ok=True)
        
        # Generate report
        editor = PPTEditor(ppt_path)
        editor.load_presentation()
        editor.replace_text_in_slide(request.slide_number, request.replacements)
        output_path = editor.save_report(output_dir)
        
        return FileResponse(
            path=output_path,
            filename=f"report_{output_path.stem}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health")
async def health_check():
    return {"status": "healthy"}